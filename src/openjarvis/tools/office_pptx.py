"""office_pptx.py — PowerPoint (.pptx) read/find/create/write for personal PC.

Actions:
1. read — extract text per slide/shape
2. find — search query across slides
3. create — new presentation with title + optional body bullets
4. write — exact text replacements or set shape text; template-safe output_path

Empty path → last opened .pptx from session_context when available.
Uses python-pptx (no PowerPoint COM).
"""

from __future__ import annotations

import json
import re
from datetime import datetime
from pathlib import Path
from typing import Any, Optional

from openjarvis.core.registry import ToolRegistry
from openjarvis.core.types import ToolResult
from openjarvis.tools._stubs import BaseTool, ToolSpec

_DEFAULT_MAX_CHARS = 50_000
_PPTX_SUFFIX = ".pptx"
_ACTIONS = ("read", "find", "create", "write")


def _fail(msg: str) -> ToolResult:
    return ToolResult(tool_name="office_pptx", content=msg, success=False)


def _ok(msg: str, **metadata: Any) -> ToolResult:
    return ToolResult(
        tool_name="office_pptx",
        content=msg,
        success=True,
        metadata=metadata or None,
    )


def _note(path: Path) -> None:
    try:
        from openjarvis.tools.session_context import note_opened

        note_opened(str(path.resolve()))
    except Exception:
        pass


def _default_output(stem: str) -> Path:
    stamp = datetime.now().strftime("%Y%m%d_%H%M%S")
    name = f"{stem}_{stamp}.pptx"
    if Path("D:/").exists():
        return Path("D:/") / name
    docs = Path.home() / "Documents"
    docs.mkdir(parents=True, exist_ok=True)
    return docs / name


def _resolve_pptx_path(
    raw: str, *, must_exist: bool = True
) -> tuple[Optional[Path], Optional[str]]:
    path_str = (raw or "").strip()
    if not path_str:
        try:
            from openjarvis.tools.session_context import get_last_opened

            last = get_last_opened()
        except Exception:
            last = None
        if last and Path(str(last)).suffix.lower() == _PPTX_SUFFIX:
            path_str = str(last)
        else:
            return None, "No path provided and no last opened .pptx in session."

    path = Path(path_str)
    if path.suffix.lower() != _PPTX_SUFFIX:
        return None, f"Not a PowerPoint .pptx file: {path_str}"

    from openjarvis.security.file_policy import is_sensitive_file

    if is_sensitive_file(path):
        return None, f"Access denied: {path_str} is a sensitive file."
    if must_exist and not path.exists():
        return None, f"File not found: {path_str}"
    return path, None


def _import_pptx() -> tuple[Any, Optional[str]]:
    try:
        from pptx import Presentation
    except ImportError:
        return None, (
            "python-pptx not installed. Install with: uv sync --extra pptx"
        )
    return Presentation, None


def _parse_pages(pages_str: str, total: int) -> list[int]:
    """1-indexed slide numbers → 0-indexed."""
    result: list[int] = []
    for part in str(pages_str).split(","):
        part = part.strip()
        if not part:
            continue
        if "-" in part:
            a, b = part.split("-", 1)
            start = max(1, int(a.strip()))
            end = min(total, int(b.strip()))
            result.extend(range(start - 1, end))
        else:
            n = int(part)
            if 1 <= n <= total:
                result.append(n - 1)
    return sorted(set(result))


def _shape_text(shape: Any) -> str:
    if not getattr(shape, "has_text_frame", False):
        return ""
    parts: list[str] = []
    for para in shape.text_frame.paragraphs:
        parts.append("".join(run.text for run in para.runs) or (para.text or ""))
    return "\n".join(p for p in parts if p is not None).strip()


def _iter_text_shapes(slide: Any) -> list[tuple[int, Any, str]]:
    out: list[tuple[int, Any, str]] = []
    for idx, shape in enumerate(slide.shapes):
        text = _shape_text(shape)
        if text:
            out.append((idx, shape, text))
    return out


def _parse_replacements(raw: Any) -> tuple[Optional[dict[str, str]], Optional[str]]:
    if raw is None:
        return None, None
    if isinstance(raw, dict):
        return {str(k): "" if v is None else str(v) for k, v in raw.items()}, None
    text = str(raw).strip()
    if not text:
        return None, None
    try:
        data = json.loads(text)
    except json.JSONDecodeError as exc:
        return None, f"replacements must be JSON object: {exc}"
    if not isinstance(data, dict):
        return None, "replacements must be a JSON object old→new."
    return {str(k): "" if v is None else str(v) for k, v in data.items()}, None


def _parse_bullets(raw: Any) -> list[str]:
    if raw is None:
        return []
    if isinstance(raw, list):
        return [str(x).strip() for x in raw if str(x).strip()]
    text = str(raw).strip()
    if not text:
        return []
    if text.startswith("["):
        try:
            data = json.loads(text)
            if isinstance(data, list):
                return [str(x).strip() for x in data if str(x).strip()]
        except json.JSONDecodeError:
            pass
    return [ln.strip(" -*\t") for ln in re.split(r"[\n;]", text) if ln.strip()]


def _replace_in_shape(shape: Any, old: str, new: str) -> int:
    if not getattr(shape, "has_text_frame", False):
        return 0
    count = 0
    for para in shape.text_frame.paragraphs:
        for run in para.runs:
            if old in run.text:
                run.text = run.text.replace(old, new)
                count += 1
        # Fallback when text lives on paragraph without runs split
        if not para.runs and para.text and old in para.text:
            para.text = para.text.replace(old, new)
            count += 1
    return count


@ToolRegistry.register("office_pptx")
class OfficePptxTool(BaseTool):
    """Personal-PC PowerPoint tool: read, find, create, write."""

    tool_id = "office_pptx"

    @property
    def spec(self) -> ToolSpec:
        return ToolSpec(
            name="office_pptx",
            description=(
                "PowerPoint .pptx на этом ПК (python-pptx). "
                "Когда использовать: прочитать слайды → read; найти текст → find; "
                "создать презентацию → create; заменить фразы → write "
                "(replacements JSON, лучше новый output_path). "
                "Когда НЕ использовать: просто открыть → open_path; "
                "PDF → office_pdf; Word → office_word; Excel → office_excel. "
                "path можно опустить, если .pptx в КОНТЕКСТЕ ФАЙЛОВ."
            ),
            parameters={
                "type": "object",
                "properties": {
                    "action": {
                        "type": "string",
                        "description": "read | find | create | write",
                    },
                    "path": {
                        "type": "string",
                        "description": "Path to .pptx (alias: file_path).",
                    },
                    "file_path": {
                        "type": "string",
                        "description": "Alias for path.",
                    },
                    "pages": {
                        "type": "string",
                        "description": "Slide range for read/find, e.g. '1-3'.",
                    },
                    "query": {
                        "type": "string",
                        "description": "Search string for action=find.",
                    },
                    "title": {
                        "type": "string",
                        "description": "Title for action=create.",
                    },
                    "body": {
                        "type": "string",
                        "description": (
                            "Body/bullets for create: JSON array or lines."
                        ),
                    },
                    "replacements": {
                        "type": "string",
                        "description": (
                            'For write: JSON object {"exact old":"new",...}.'
                        ),
                    },
                    "slide": {
                        "type": "integer",
                        "description": "1-based slide for write shape text.",
                    },
                    "shape_index": {
                        "type": "integer",
                        "description": "0-based shape index for write text=.",
                    },
                    "text": {
                        "type": "string",
                        "description": "New text when setting a shape via write.",
                    },
                    "output_path": {
                        "type": "string",
                        "description": "Output .pptx for create/write.",
                    },
                    "max_chars": {
                        "type": "integer",
                        "description": "Max characters to return.",
                    },
                    "overwrite": {
                        "type": "boolean",
                        "description": "Allow saving write back to path.",
                    },
                },
                "required": ["action"],
            },
            category="office",
            required_capabilities=["file:read"],
        )

    def execute(self, **params: Any) -> ToolResult:
        action = str(params.get("action") or "").strip().lower()
        if action not in _ACTIONS:
            return _fail(
                f"Unknown action {action!r}. Use: {', '.join(_ACTIONS)}"
            )
        if action == "create":
            return self._create(params)
        if action == "write":
            return self._write(params)
        if action == "find":
            return self._find(params)
        return self._read(params)

    def _read(self, params: Any) -> ToolResult:
        path, err = _resolve_pptx_path(
            str(params.get("path") or params.get("file_path") or "")
        )
        if err or path is None:
            return _fail(err or "No path.")
        Presentation, ierr = _import_pptx()
        if ierr:
            return _fail(ierr)

        max_chars = int(params.get("max_chars") or _DEFAULT_MAX_CHARS)
        try:
            prs = Presentation(str(path))
            total = len(prs.slides)
            pages_param = params.get("pages")
            if pages_param:
                indices = _parse_pages(str(pages_param), total)
            else:
                indices = list(range(total))

            blocks: list[str] = []
            for idx in indices:
                slide = prs.slides[idx]
                blocks.append(f"--- slide {idx + 1} ---")
                shapes = _iter_text_shapes(slide)
                if not shapes:
                    blocks.append("(no text shapes)")
                    continue
                for sidx, _shape, text in shapes:
                    blocks.append(f"[shape {sidx}]\n{text}")
            body = "\n\n".join(blocks)
            if len(body) > max_chars:
                body = body[:max_chars] + "\n\n[Content truncated]"
        except Exception as exc:
            return _fail(f"Read failed: {exc}")

        _note(path)
        return _ok(
            body or f"No slides in {path}.",
            file_path=str(path.resolve()),
            slide_count=total,
        )

    def _find(self, params: Any) -> ToolResult:
        query = str(params.get("query") or "").strip()
        if not query:
            return _fail("query is required for action=find.")
        path, err = _resolve_pptx_path(
            str(params.get("path") or params.get("file_path") or "")
        )
        if err or path is None:
            return _fail(err or "No path.")
        Presentation, ierr = _import_pptx()
        if ierr:
            return _fail(ierr)

        max_chars = int(params.get("max_chars") or _DEFAULT_MAX_CHARS)
        q = query.lower()
        hits: list[str] = []
        try:
            prs = Presentation(str(path))
            total = len(prs.slides)
            pages_param = params.get("pages")
            indices = (
                _parse_pages(str(pages_param), total)
                if pages_param
                else list(range(total))
            )
            for idx in indices:
                for sidx, _shape, text in _iter_text_shapes(prs.slides[idx]):
                    if q in text.lower():
                        for line in text.splitlines():
                            if q in line.lower():
                                hits.append(
                                    f"slide {idx + 1} shape {sidx}: {line.strip()}"
                                )
        except Exception as exc:
            return _fail(f"Find failed: {exc}")

        _note(path)
        if not hits:
            return _ok(
                f"No matches for {query!r} in {path}.",
                file_path=str(path.resolve()),
                match_count=0,
            )
        body = f"Matches for {query!r} ({len(hits)}):\n" + "\n".join(hits)
        if len(body) > max_chars:
            body = body[:max_chars] + "\n\n[Content truncated]"
        return _ok(
            body,
            file_path=str(path.resolve()),
            match_count=len(hits),
        )

    def _create(self, params: Any) -> ToolResult:
        Presentation, ierr = _import_pptx()
        if ierr:
            return _fail(ierr)

        out_raw = str(
            params.get("output_path")
            or params.get("path")
            or params.get("file_path")
            or ""
        ).strip()
        out_path = Path(out_raw) if out_raw else _default_output("presentation")
        if out_path.suffix.lower() != _PPTX_SUFFIX:
            out_path = out_path.with_suffix(_PPTX_SUFFIX)
        if out_path.exists() and not bool(params.get("overwrite")):
            return _fail(
                f"File exists: {out_path}. Pass overwrite=true or another path."
            )
        out_path.parent.mkdir(parents=True, exist_ok=True)

        title = str(params.get("title") or "Presentation").strip() or "Presentation"
        bullets = _parse_bullets(params.get("body"))

        try:
            prs = Presentation()
            layout = prs.slide_layouts[1]  # title and content
            slide = prs.slides.add_slide(layout)
            slide.shapes.title.text = title
            if len(slide.placeholders) > 1 and bullets:
                body = slide.placeholders[1]
                tf = body.text_frame
                tf.clear()
                for i, line in enumerate(bullets):
                    if i == 0:
                        tf.text = line
                    else:
                        p = tf.add_paragraph()
                        p.text = line
                        p.level = 0
            elif not bullets and len(slide.placeholders) > 1:
                slide.placeholders[1].text = ""
            prs.save(str(out_path))
        except Exception as exc:
            return _fail(f"Create failed: {exc}")

        _note(out_path)
        return _ok(
            f"Created presentation: {out_path.resolve()} "
            f"(title={title!r}, bullets={len(bullets)})",
            file_path=str(out_path.resolve()),
        )

    def _write(self, params: Any) -> ToolResult:
        path, err = _resolve_pptx_path(
            str(params.get("path") or params.get("file_path") or "")
        )
        if err or path is None:
            return _fail(err or "No path.")
        Presentation, ierr = _import_pptx()
        if ierr:
            return _fail(ierr)

        replacements, rerr = _parse_replacements(params.get("replacements"))
        if rerr:
            return _fail(rerr)

        slide_raw = params.get("slide")
        shape_index = params.get("shape_index")
        new_text = params.get("text")
        set_shape = (
            slide_raw is not None
            and shape_index is not None
            and new_text is not None
        )
        if not replacements and not set_shape:
            return _fail(
                "write needs replacements JSON and/or "
                "slide+shape_index+text."
            )

        overwrite = bool(params.get("overwrite"))
        out_raw = str(params.get("output_path") or "").strip()
        if out_raw:
            out_path = Path(out_raw)
        elif overwrite:
            out_path = path
        else:
            out_path = _default_output(path.stem + "_edit")

        if out_path.suffix.lower() != _PPTX_SUFFIX:
            out_path = out_path.with_suffix(_PPTX_SUFFIX)

        if out_path.resolve() == path.resolve() and not overwrite:
            return _fail(
                "Refusing to overwrite source. Pass overwrite=true "
                "or a different output_path."
            )
        out_path.parent.mkdir(parents=True, exist_ok=True)

        try:
            prs = Presentation(str(path))
            changed = 0
            if set_shape:
                slide_no = int(slide_raw)
                sidx = int(shape_index)
                if slide_no < 1 or slide_no > len(prs.slides):
                    return _fail(f"slide {slide_no} out of range.")
                slide = prs.slides[slide_no - 1]
                if sidx < 0 or sidx >= len(slide.shapes):
                    return _fail(f"shape_index {sidx} out of range.")
                shape = slide.shapes[sidx]
                if not getattr(shape, "has_text_frame", False):
                    return _fail(f"shape {sidx} has no text frame.")
                shape.text_frame.text = str(new_text)
                changed += 1

            if replacements:
                # Longest keys first (same idea as Word clone).
                keys = sorted(replacements.keys(), key=len, reverse=True)
                for slide in prs.slides:
                    for shape in slide.shapes:
                        for old in keys:
                            changed += _replace_in_shape(
                                shape, old, replacements[old]
                            )

            prs.save(str(out_path))
        except Exception as exc:
            return _fail(f"Write failed: {exc}")

        _note(out_path)
        return _ok(
            f"Wrote presentation ({changed} change(s)) → {out_path.resolve()}",
            file_path=str(out_path.resolve()),
            source=str(path.resolve()),
            changes=changed,
        )


__all__ = ["OfficePptxTool"]
